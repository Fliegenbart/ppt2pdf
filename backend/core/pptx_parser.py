"""
PPTX Parser - Extracts all elements from PowerPoint presentations.
"""
import base64
import io
import uuid
from typing import Optional
from pptx import Presentation as PPTXPresentation
from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.shapes.base import BaseShape
from pptx.shapes.picture import Picture
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.shapes.autoshape import Shape
from pptx.table import Table
from PIL import Image

from api.models import (
    Presentation, Slide, SlideElement, ElementType, ContentType,
    BoundingBox, TextParagraph, TextRun, TextStyle,
    TableData, TableCell, ChartData
)


class PPTXParser:
    """Parses PowerPoint files and extracts structured content."""

    def __init__(self):
        self.element_counter = 0

    def parse(self, file_path: str) -> Presentation:
        """Parse a PPTX file and return structured presentation data."""
        pptx = PPTXPresentation(file_path)

        # Extract presentation metadata
        presentation = Presentation(
            filename=file_path.split("/")[-1],
            title=self._get_core_property(pptx, "title"),
            author=self._get_core_property(pptx, "author"),
            slides=[],
        )

        # Parse each slide
        for slide_num, slide in enumerate(pptx.slides, start=1):
            parsed_slide = self._parse_slide(slide, slide_num)
            presentation.slides.append(parsed_slide)

        return presentation

    def _get_core_property(self, pptx: PPTXPresentation, prop: str) -> Optional[str]:
        """Safely get a core property from the presentation."""
        try:
            return getattr(pptx.core_properties, prop, None)
        except Exception:
            return None

    def _parse_slide(self, slide, slide_number: int) -> Slide:
        """Parse a single slide and extract all elements."""
        parsed_slide = Slide(
            slide_number=slide_number,
            elements=[],
        )

        # Get slide title
        if slide.shapes.title:
            parsed_slide.title = slide.shapes.title.text

        # Get speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            parsed_slide.speaker_notes = slide.notes_slide.notes_text_frame.text

        # Get layout name
        try:
            parsed_slide.layout_name = slide.slide_layout.name
        except Exception:
            pass

        # Get background color
        parsed_slide.background_color = self._get_background_color(slide)

        # Parse all shapes
        for shape in slide.shapes:
            elements = self._parse_shape(shape, slide_number)
            parsed_slide.elements.extend(elements)

        # Assign initial reading order based on position (top-to-bottom, left-to-right)
        self._assign_initial_reading_order(parsed_slide.elements)

        return parsed_slide

    def _get_background_color(self, slide) -> Optional[str]:
        """Extract background color from slide."""
        try:
            fill = slide.background.fill
            if fill.type is not None:
                if hasattr(fill, 'fore_color') and fill.fore_color:
                    return self._color_to_hex(fill.fore_color)
        except Exception:
            pass
        return "#FFFFFF"  # Default white

    def _parse_shape(self, shape: BaseShape, slide_number: int) -> list[SlideElement]:
        """Parse a shape and return one or more elements."""
        elements = []

        # Get bounding box
        bounds = BoundingBox(
            x=float(shape.left) if shape.left else 0,
            y=float(shape.top) if shape.top else 0,
            width=float(shape.width) if shape.width else 0,
            height=float(shape.height) if shape.height else 0,
        )

        element_id = f"slide{slide_number}_elem{self._next_id()}"

        # Handle different shape types
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            element = self._parse_picture(shape, element_id, bounds)
            if element:
                elements.append(element)

        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            element = self._parse_table(shape, element_id, bounds)
            if element:
                elements.append(element)

        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            element = self._parse_chart(shape, element_id, bounds)
            if element:
                elements.append(element)

        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # Recursively parse group shapes
            group_elements = self._parse_group(shape, slide_number)
            elements.extend(group_elements)

        elif hasattr(shape, "text_frame"):
            # Text-containing shapes
            element = self._parse_text_shape(shape, element_id, bounds)
            if element:
                elements.append(element)

        elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            # Handle placeholders
            element = self._parse_placeholder(shape, element_id, bounds)
            if element:
                elements.append(element)

        return elements

    def _parse_picture(self, shape: Picture, element_id: str, bounds: BoundingBox) -> Optional[SlideElement]:
        """Parse an image/picture shape."""
        try:
            # Extract image data
            image = shape.image
            image_bytes = image.blob
            image_base64 = base64.b64encode(image_bytes).decode('utf-8')

            # Determine content type based on image analysis
            content_type = self._classify_image(image_bytes)

            # Check for existing alt text
            alt_text = None
            try:
                if hasattr(shape, '_element'):
                    desc_elem = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}cNvPr')
                    if desc_elem is not None:
                        alt_text = desc_elem.get('descr')
            except Exception:
                pass

            return SlideElement(
                id=element_id,
                element_type=ElementType.IMAGE,
                bounds=bounds,
                image_base64=image_base64,
                content_type=content_type,
                alt_text=alt_text,
                is_decorative=False,
            )
        except Exception as e:
            print(f"Error parsing picture: {e}")
            return None

    def _parse_text_shape(self, shape, element_id: str, bounds: BoundingBox) -> Optional[SlideElement]:
        """Parse a shape with text content."""
        if not hasattr(shape, 'text_frame'):
            return None

        text_frame = shape.text_frame
        if not text_frame.paragraphs:
            return None

        paragraphs = []
        has_content = False

        for para in text_frame.paragraphs:
            runs = []
            for run in para.runs:
                if run.text.strip():
                    has_content = True
                    style = self._extract_text_style(run)
                    runs.append(TextRun(text=run.text, style=style))

            if runs:
                # Check for bullet
                is_bullet = para.level > 0 or self._has_bullet(para)
                bullet_char = self._get_bullet_char(para) if is_bullet else None

                paragraphs.append(TextParagraph(
                    runs=runs,
                    level=para.level,
                    is_bullet=is_bullet,
                    bullet_char=bullet_char,
                ))

        if not has_content:
            return None

        # Determine if this is a heading based on style
        heading_level = self._detect_heading_level(shape, paragraphs)

        return SlideElement(
            id=element_id,
            element_type=ElementType.TEXT,
            bounds=bounds,
            paragraphs=paragraphs,
            heading_level=heading_level,
        )

    def _parse_table(self, shape, element_id: str, bounds: BoundingBox) -> Optional[SlideElement]:
        """Parse a table shape."""
        try:
            if not isinstance(shape, GraphicFrame):
                return None

            table = shape.table
            rows = []
            has_header_row = False
            has_header_column = False

            for row_idx, row in enumerate(table.rows):
                row_cells = []
                for col_idx, cell in enumerate(row.cells):
                    # Check if this looks like a header
                    is_header = row_idx == 0 or col_idx == 0
                    if row_idx == 0:
                        has_header_row = True
                    if col_idx == 0 and self._is_header_cell(cell):
                        has_header_column = True

                    style = self._extract_cell_style(cell)
                    row_cells.append(TableCell(
                        text=cell.text,
                        is_header=is_header,
                        style=style,
                    ))
                rows.append(row_cells)

            table_data = TableData(
                rows=rows,
                has_header_row=has_header_row,
                has_header_column=has_header_column,
            )

            return SlideElement(
                id=element_id,
                element_type=ElementType.TABLE,
                bounds=bounds,
                table_data=table_data,
            )
        except Exception as e:
            print(f"Error parsing table: {e}")
            return None

    def _parse_chart(self, shape, element_id: str, bounds: BoundingBox) -> Optional[SlideElement]:
        """Parse a chart shape."""
        try:
            if not isinstance(shape, GraphicFrame):
                return None

            chart = shape.chart
            chart_data = ChartData(
                chart_type=str(chart.chart_type) if hasattr(chart, 'chart_type') else "unknown",
                title=chart.chart_title.text_frame.text if chart.has_title else None,
            )

            # Extract series data if available
            try:
                for series in chart.series:
                    series_data = {
                        "name": series.name,
                        "values": list(series.values) if hasattr(series, 'values') else [],
                    }
                    chart_data.series.append(series_data)

                # Extract categories
                if hasattr(chart, 'plots') and chart.plots:
                    for plot in chart.plots:
                        if hasattr(plot, 'categories'):
                            chart_data.categories = list(plot.categories)
                            break
            except Exception:
                pass

            return SlideElement(
                id=element_id,
                element_type=ElementType.CHART,
                bounds=bounds,
                chart_data=chart_data,
            )
        except Exception as e:
            print(f"Error parsing chart: {e}")
            return None

    def _parse_group(self, shape: GroupShape, slide_number: int) -> list[SlideElement]:
        """Parse a group of shapes."""
        elements = []
        for child_shape in shape.shapes:
            child_elements = self._parse_shape(child_shape, slide_number)
            elements.extend(child_elements)
        return elements

    def _parse_placeholder(self, shape, element_id: str, bounds: BoundingBox) -> Optional[SlideElement]:
        """Parse a placeholder shape."""
        if hasattr(shape, 'text_frame') and shape.text_frame:
            return self._parse_text_shape(shape, element_id, bounds)
        return None

    def _extract_text_style(self, run) -> TextStyle:
        """Extract text styling from a run."""
        style = TextStyle()
        try:
            font = run.font
            style.font_name = font.name
            style.font_size = font.size.pt if font.size else None
            style.bold = font.bold or False
            style.italic = font.italic or False
            style.underline = font.underline or False
            if font.color and font.color.rgb:
                style.color = self._rgb_to_hex(font.color.rgb)
        except Exception:
            pass
        return style

    def _extract_cell_style(self, cell) -> TextStyle:
        """Extract style from a table cell."""
        style = TextStyle()
        try:
            if cell.text_frame and cell.text_frame.paragraphs:
                para = cell.text_frame.paragraphs[0]
                if para.runs:
                    return self._extract_text_style(para.runs[0])
        except Exception:
            pass
        return style

    def _classify_image(self, image_bytes: bytes) -> ContentType:
        """Classify the type of image content."""
        try:
            img = Image.open(io.BytesIO(image_bytes))
            width, height = img.size

            # Simple heuristics - AI will do better classification later
            aspect_ratio = width / height if height > 0 else 1

            # Very small images are likely icons
            if width < 100 and height < 100:
                return ContentType.ICON

            # Very wide or tall images might be decorative
            if aspect_ratio > 4 or aspect_ratio < 0.25:
                return ContentType.DECORATIVE

            return ContentType.UNKNOWN
        except Exception:
            return ContentType.UNKNOWN

    def _detect_heading_level(self, shape, paragraphs: list[TextParagraph]) -> Optional[int]:
        """Detect if this text element is a heading and what level."""
        try:
            # Check if it's a title placeholder
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type is not None:
                    ph_name = str(ph_type).lower()
                    if 'title' in ph_name:
                        return 1
                    if 'subtitle' in ph_name:
                        return 2

            # Check font size heuristics
            if paragraphs and paragraphs[0].runs:
                font_size = paragraphs[0].runs[0].style.font_size
                if font_size:
                    if font_size >= 32:
                        return 1
                    elif font_size >= 24:
                        return 2
                    elif font_size >= 18:
                        return 3
        except Exception:
            pass
        return None

    def _has_bullet(self, para) -> bool:
        """Check if paragraph has a bullet."""
        try:
            return para._pPr is not None and para._pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buChar') is not None
        except Exception:
            return False

    def _get_bullet_char(self, para) -> Optional[str]:
        """Get the bullet character if any."""
        try:
            bu = para._pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
            if bu is not None:
                return bu.get('char')
        except Exception:
            pass
        return "•"

    def _is_header_cell(self, cell) -> bool:
        """Check if a cell appears to be a header."""
        try:
            if cell.text_frame and cell.text_frame.paragraphs:
                para = cell.text_frame.paragraphs[0]
                if para.runs:
                    return para.runs[0].font.bold or False
        except Exception:
            pass
        return False

    def _color_to_hex(self, color) -> Optional[str]:
        """Convert a color object to hex string."""
        try:
            if hasattr(color, 'rgb') and color.rgb:
                return self._rgb_to_hex(color.rgb)
        except Exception:
            pass
        return None

    def _rgb_to_hex(self, rgb: RGBColor) -> str:
        """Convert RGBColor to hex string."""
        return f"#{rgb}"

    def _assign_initial_reading_order(self, elements: list[SlideElement]):
        """Assign initial reading order based on position."""
        # Sort by y position (top to bottom), then x (left to right)
        sorted_elements = sorted(
            elements,
            key=lambda e: (e.bounds.y, e.bounds.x)
        )
        for order, element in enumerate(sorted_elements):
            element.reading_order = order

    def _next_id(self) -> int:
        """Generate next element ID."""
        self.element_counter += 1
        return self.element_counter
