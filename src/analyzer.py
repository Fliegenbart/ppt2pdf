"""
PPTX-Analyse für Barrierefreiheit.
"""

import logging
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .alt_text import CachedAltTextGenerator, LocalAltTextGenerator
from .config import Config
from .models import AnalysisResult, SlideContent, SlideImage

logger = logging.getLogger(__name__)


class PPTXAnalyzer:
    """Analysiert PowerPoint-Präsentationen für Barrierefreiheit."""

    def __init__(self, config: Config):
        self.config = config

        # Alt-Text Generator (mit oder ohne Cache)
        self.alt_generator: CachedAltTextGenerator | LocalAltTextGenerator
        if config.use_cache:
            self.alt_generator = CachedAltTextGenerator(config)
        else:
            self.alt_generator = LocalAltTextGenerator(config)

    def analyze(self, pptx_path: str | Path) -> AnalysisResult:
        """
        Analysiert eine PPTX-Datei komplett.

        Args:
            pptx_path: Pfad zur PowerPoint-Datei

        Returns:
            AnalysisResult mit allen extrahierten Informationen
        """
        pptx_path = Path(pptx_path)

        if not pptx_path.exists():
            raise FileNotFoundError(f"PPTX nicht gefunden: {pptx_path}")

        prs = Presentation(str(pptx_path))

        result = AnalysisResult(
            slides=[],
            total_slides=len(prs.slides),
        )

        # Präsentationstitel für Kontext
        result.presentation_title = self._extract_presentation_title(prs)

        # Slides analysieren
        for slide_idx, slide in enumerate(prs.slides):
            slide_content = self._analyze_slide(
                slide,
                slide_idx,
                result.presentation_title
            )
            result.slides.append(slide_content)

            # Statistiken aktualisieren
            result.images_total += slide_content.image_count
            result.images_without_alt += len(slide_content.images_without_alt)

            if slide_content.has_table:
                result.tables += 1
            if slide_content.has_chart:
                result.charts += 1

        # Zähle generierte Alt-Texte
        for slide in result.slides:
            for img in slide.images:
                if img.generated_alt_text:
                    result.images_alt_generated += 1

        return result

    def _extract_presentation_title(self, prs: Presentation) -> str:
        """Extrahiert den Präsentationstitel aus der ersten Folie."""
        if not prs.slides:
            return ""

        first_slide = prs.slides[0]

        # Versuche Titel-Shape zu finden
        if first_slide.shapes.title and first_slide.shapes.title.text:
            return str(first_slide.shapes.title.text).strip()[:100]

        # Fallback: Erster Text auf der Folie
        for shape in first_slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                return str(shape.text).strip()[:100]

        return ""

    def _analyze_slide(
        self,
        slide,
        slide_idx: int,
        presentation_title: str
    ) -> SlideContent:
        """Analysiert eine einzelne Folie."""

        slide_content = SlideContent(index=slide_idx)

        # Titel extrahieren
        if slide.shapes.title:
            slide_content.title = slide.shapes.title.text or ""

        for shape_idx, shape in enumerate(slide.shapes):
            # Text extrahieren
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text and text != slide_content.title:
                        slide_content.text_content.append(text)

            # Bilder verarbeiten
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_image = self._process_image(
                    shape,
                    slide_idx,
                    shape_idx,
                    presentation_title,
                    slide_content.title,
                )
                if slide_image:
                    slide_content.images.append(slide_image)

            # Tabellen
            if shape.has_table:
                slide_content.has_table = True

            # Charts
            if shape.has_chart:
                slide_content.has_chart = True

        # Speaker Notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame
            if notes:
                slide_content.speaker_notes = notes.text

        return slide_content

    def _process_image(
        self,
        shape,
        slide_idx: int,
        shape_idx: int,
        presentation_title: str,
        slide_title: str,
    ) -> SlideImage | None:
        """Verarbeitet ein Bild-Shape und generiert ggf. Alt-Text."""

        try:
            image = shape.image

            # Aktuellen Alt-Text extrahieren (aus Shape-Name als Fallback)
            current_alt = ""
            if hasattr(shape, "name"):
                current_alt = shape.name or ""

            # Prüfen ob Alt-Text in XML vorhanden (descr-Attribut)
            try:
                element = shape._element
                nvPicPr = element.find(".//{http://schemas.openxmlformats.org/presentationml/2006/main}nvPicPr")
                if nvPicPr is not None:
                    cNvPr = nvPicPr.find("{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr")
                    if cNvPr is not None and cNvPr.get("descr"):
                        current_alt = cNvPr.get("descr")
            except Exception:
                pass

            slide_image = SlideImage(
                slide_index=slide_idx,
                shape_index=shape_idx,
                image_bytes=image.blob,
                current_alt_text=current_alt,
                content_type=image.content_type,
            )

            # Alt-Text generieren wenn nötig
            needs_generation = (
                not current_alt or not self.config.skip_existing_alt_texts
            )

            if needs_generation and self.alt_generator.is_available():
                context = f"{presentation_title} - {slide_title}".strip(" -")

                logger.info(f"  🖼️  Folie {slide_idx + 1}: Generiere Alt-Text...")

                generated = self.alt_generator.generate(image.blob, context)

                if generated:
                    slide_image.generated_alt_text = generated
                    logger.info(f"    ✓ \"{generated[:50]}...\"")

            return slide_image

        except Exception as e:
            logger.warning(f"Bild-Extraktion fehlgeschlagen: {e}")
            return None
